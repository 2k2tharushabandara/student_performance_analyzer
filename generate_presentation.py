from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

root = r"f:\ML_Assignment_Student_Performance_Analyzer_Group_07"
output_path = os.path.join(root, "Student_Performance_Presentation_Group_07.pptx")

prs = Presentation()


def set_title(slide, text):
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.size = Pt(40)


def add_bullets(slide, bullets, level=0):
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, item in enumerate(bullets):
        p = body.add_paragraph() if i > 0 else body.paragraphs[0]
        p.text = item
        p.level = level
        p.font.size = Pt(20)


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)
    slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(20)
    return slide


def add_bullet_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, title)
    add_bullets(slide, bullets)
    return slide


def add_image_slide(title, image_path, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    set_title(slide, title)

    left = Inches(0.6)
    top = Inches(1.4)
    width = Inches(9.0)
    try:
        slide.shapes.add_picture(image_path, left, top, width=width)
    except FileNotFoundError:
        tx = slide.shapes.add_textbox(left, top, width, Inches(4.5))
        tx.text_frame.text = f"[Missing image: {os.path.basename(image_path)}]"

    if caption:
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(6.7), Inches(9.0), Inches(0.5))
        tf = tx.text_frame
        tf.text = caption
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return slide


# Title slide
add_title_slide(
    "Student Performance Analyzer — Grade Classification",
    "Group 07 • May 25, 2026"
)

# Agenda
add_bullet_slide("Agenda", [
    "Problem & objectives",
    "Dataset & features",
    "Exploratory analysis",
    "Preprocessing & balancing",
    "Model comparison & selection",
    "Evaluation results",
    "Deployment & next steps",
])

# Problem & objective
add_bullet_slide("Problem & Objective", [
    "Goal: predict student grade category (A, B, C, D, Fail).",
    "Use student attributes (study habits, attendance, GPA, etc.) to classify outcomes.",
    "Emphasis on interpretability + deployable inference (API/UI).",
])

# Dataset overview
add_bullet_slide("Dataset Overview", [
    "Source file: student_performance_grade.csv",
    "Labeled rows (classification): 1,671",
    "Features: 18 (numeric + categorical + binary)",
    "Target: Grade (A, B, C, D, Fail) — imbalanced distribution",
])

# EDA slides
add_image_slide(
    "EDA — Grade Distribution",
    os.path.join(root, "eda_grade_dist.png"),
    "Class imbalance: A dominates; D/Fail are rare"
)

add_image_slide(
    "EDA — Numerical Feature Distributions",
    os.path.join(root, "eda_num_distributions.png"),
    "Key numeric patterns (hours studied, attendance, GPA, etc.)"
)

add_image_slide(
    "EDA — Categorical Feature Boxplots",
    os.path.join(root, "eda_categorical_boxplots.png"),
    "Grade variations across categorical attributes"
)

add_image_slide(
    "EDA — Correlation Heatmap",
    os.path.join(root, "eda_correlation_heatmap.png"),
    "Correlations among numeric variables and final score"
)

# Preprocessing
add_bullet_slide("Preprocessing & Balancing", [
    "Encoding: ordinal + one-hot + binary mapping",
    "Scaling: standardization for numeric features",
    "Train/test split with stratification",
    "Class imbalance handled via SMOTE",
    "Before SMOTE: A=921, B=477, C=219, D=46, Fail=8",
    "After SMOTE: all classes balanced to 921 each",
])

# Model comparison
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_title(slide, "Model Comparison (CV F1 — weighted)")

rows, cols = 5, 3
left = Inches(0.8)
top = Inches(1.6)
width = Inches(8.4)
height = Inches(2.0)

table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
_table = table_shape.table

_table.cell(0, 0).text = "Model"
_table.cell(0, 1).text = "CV F1 (mean)"
_table.cell(0, 2).text = "Std"

models = [
    ("Logistic Regression", "0.8300", "0.0207"),
    ("K-Nearest Neighbors", "0.8751", "0.0333"),
    ("Random Forest", "0.9237", "0.0300"),
    ("Gradient Boosting", "0.8719", "0.0320"),
]
for i, (m, mean, std) in enumerate(models, start=1):
    _table.cell(i, 0).text = m
    _table.cell(i, 1).text = mean
    _table.cell(i, 2).text = std

# Tune/selection
add_bullet_slide("Best Model — Random Forest (Tuned)", [
    "Best params: n_estimators=200, max_depth=None, min_samples_split=2",
    "Best CV F1 (weighted): 0.9264",
    "Selected for balance of performance + interpretability",
])

# Test set results
add_bullet_slide("Test Set Performance", [
    "Accuracy: 0.73",
    "Weighted F1: 0.72",
    "Macro F1: 0.49 (minority classes remain challenging)",
    "Low recall for D and Fail reflects rare class scarcity",
])

# Confusion matrix
add_image_slide(
    "Confusion Matrix — Grade Classification",
    os.path.join(root, "eval_clf_confusion_matrix.png"),
    "Most confusion occurs among adjacent grades; minority classes hardest"
)

# Feature importance
add_image_slide(
    "Top Feature Importances (Random Forest)",
    os.path.join(root, "eval_clf_feature_importance.png"),
    "Most influential features include GPA, attendance, study time, and anxiety"
)

# Deployment
add_bullet_slide("Deployment", [
    "Model saved: model_classification.pkl + preprocessor_classification.pkl",
    "Gradio app UI: app.py",
    "Flask API endpoint: app_flask.py (/predict)",
    "Supports JSON input → predicted grade + probabilities",
])

# Sample prediction
add_bullet_slide("Sample Prediction", [
    "Example input (Age=20, Study Method=Hybrid, Attendance=85%, GPA=3.2, …)",
    "Predicted Grade: A",
    "Probabilities: A=0.980, B=0.020, C=0.000, D=0.000, Fail=0.000",
])

# Limitations & next steps
add_bullet_slide("Limitations & Next Steps", [
    "Minority classes (D/Fail) underrepresented → low recall",
    "Collect more data for rare classes; explore cost-sensitive learning",
    "Calibrate probabilities; evaluate with macro/weighted metrics",
    "Add new features (behavioral, longitudinal trends) for robustness",
])

# Q&A
add_bullet_slide("Q&A", ["Thank you! Questions?"])

prs.save(output_path)
print(output_path)
